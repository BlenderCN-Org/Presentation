import json
from pptx import Presentation
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.shapes.autoshape import Shape
from pptx.shapes.picture import Picture
from pptx.text.text import Font, TextFrame
from pptx.dml.color import _NoneColor
from pptx.chart.chart import Chart
from pptx.chart.series import BarSeries, LineSeries

from pptx.shapes.graphfrm import GraphicFrame
from pptx.util import Length, Centipoints, Cm, Emu, Inches, Mm, Pt, Px
import pptx
import inspect
prs = Presentation('existing.pptx')

def iter_rPrs(txBody):
    for p in txBody.p_lst:
        for elm in p.content_children:
            yield elm.get_or_add_rPr()
        # generate a:endParaRPr for each <a:p> element
        yield p.get_or_add_endParaRPr()
def captureColor(color):
    result = {}
    print(color)
    print(color.type)
    print(color._color)
    if not isinstance(color._color, _NoneColor):
        result["rgb"] = color.rgb
    return result    
def captureFill(fill):
    result = {}
    print(fill.type)
    if fill.type != None:
        print(fill.fore_color)
        result["fore_color"] = captureColor(fill.fore_color)
    return result

output_slides = []
print("//////////////////////////////////////////////////////////// slide")
for slide in prs.slides:
    print("------------------- slide")
    output_slide = []
    output_slides.append({"slide": output_slide})
    for shape in slide.shapes:
        print("%%%%%%%%%%%%%%%%%% - shape")
        print(shape)
        panel = {}
        output_slide.append(panel)
        if isinstance(shape, GraphicFrame):
            if shape.has_chart:
                chart_part = shape.chart_part
                print(chart_part)
                chart = chart_part.chart;
                if(chart.has_legend):
                    print("chart has a legend")
                    legend = chart.legend
                    legendInfo = {"font": legend.font.name}
                    panel["legend"] = legendInfo
                panel["chart_style"] = chart.chart_style   
                c_series = []
                panel["series"] = c_series
                for serie in chart.series:
                    c_serie = {}
                    c_series.append(c_serie)
                    c_serie["smooth"] = False
                    if isinstance(serie, LineSeries):
                        c_serie["smooth"] = serie.smooth
                        c_serie["type"] = "lineseries"
                        
                    c_serie["fill"] = None
                    if isinstance(serie, BarSeries):
                        fill_format = {}
                        c_serie["fill"] = fill_format
                        c_serie["type"] = "barseries"
                        if serie.line != None:
                            fill_format["color"] = captureColor(serie.line.color)
                            fill_format["fill"] = captureFill(serie.line.fill)
                            fill_format["width"] = serie.line.width
                            
                        # c_serie["fill"] = serie.fill
                    c_serie["values"] = serie.values
                    c_serie["name"] = serie.name
                    c_serie["index"] = serie.index
                    
            if shape.has_table:
                print("has a table")
        if isinstance(shape, Picture):
            panel["filename"] = shape.image.filename
            panel["fileext"] = shape.image.ext
            panel["contenttype"] = shape.image.content_type
            # print(shape.image.blob)
            # with open(shape.name + "." + shape.image.ext, "wb") as blob_file:
            #     blob_file.write(shape.image.blob)
            print(shape.image)
        if hasattr(shape, "text_frame"):
            panel["auto_size"] = "{}".format(shape.text_frame.auto_size)
            if isinstance(shape.text_frame, TextFrame):
                _paragraphs = []
                panel["paragraphs"] = _paragraphs
                for paragraph in shape.text_frame.paragraphs:
                    _paragraph = {}
                    _paragraphs.append(_paragraph)
                    fff = paragraph.font
                    try:
                        _paragraph["font"] = fff.name
                    except Exception:
                        _paragraph["font"] = null
                    if fff.size != None:
                        _paragraph["size"] = Length(fff.size).pt
                    _paragraph["fill"] = captureFill(fff.fill)
            for txBody in shape._element:
                try:
                    if hasattr(txBody, "p_lst"):
                        for rPr in iter_rPrs(txBody):
                            fon = Font(rPr)
                except:
                    print("error")
            panel["marginBottom"] = Length(shape.text_frame.margin_bottom).pt
            panel["marginTop"] = Length(shape.text_frame.margin_top).pt
            panel["marginLeft"] = Length(shape.text_frame.margin_left).pt
            panel["marginRight"] = Length(shape.text_frame.margin_right).pt
            
        print(shape.part)
        if hasattr(shape, "shape_type"):
            panel["shape_type"] = "{}".format(shape.shape_type)
        if hasattr(shape, "text"):
            print("text {}".format(shape.text))
            panel["text"] = shape.text
        if hasattr(shape, "height"):
            print("height {}".format(Length(shape.height).pt))
            panel["height"] = Length(shape.height).pt
        if hasattr(shape, "width"):
            print("width {}".format(Length(shape.width).pt))
            panel["width"] = Length(shape.width).pt
        if hasattr(shape, "left"):
            print("left {}".format(Length(shape.left).pt))
            panel["left"] = Length(shape.left).pt
        if hasattr(shape, "top"):
            print("top {}".format(Length(shape.top).pt))
            panel["top"] = Length(shape.top).pt
        if hasattr(shape, "rotation"):
            print("rotation {}".format(shape.rotation))
            panel["rotation"] = shape.rotation
        if hasattr(shape.part, "name"):
            print("name {}".format(shape.part.name))
            panel["name"] = shape.part.name or shape.name

# print(json.dumps(output_slides, sort_keys=True, indent=4, separators=(',', ': ')))
# prs.save('test.pptx')
with open("Output.json", "w") as text_file:
    text_file.write(json.dumps(output_slides, sort_keys=True, indent=4, separators=(',', ': ')))
